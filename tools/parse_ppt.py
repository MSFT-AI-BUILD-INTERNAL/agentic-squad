#!/usr/bin/env python3
"""PowerPoint Parser Tool for Copilot CLI.

Parses .pptx files and outputs structured data (JSON or Markdown)
that Copilot CLI agents can consume and reason about.

Usage:
    python tools/parse_ppt.py <file.pptx> [options]

Options:
    --slides           List all slide numbers and titles
    --slide N          Parse specific slide number (1-based)
    --format FORMAT    Output format: json, markdown (default: json)
    --images           Include image metadata (filename, size, position)
    --tables           Extract tables only
    --notes            Include speaker notes
    --summary          Output summary (slide count, structure overview)
    --max-slides N     Maximum slides to output (default: unlimited)

Examples:
    python tools/parse_ppt.py presentation.pptx
    python tools/parse_ppt.py presentation.pptx --slides
    python tools/parse_ppt.py presentation.pptx --slide 3 --format markdown
    python tools/parse_ppt.py presentation.pptx --tables --format json
    python tools/parse_ppt.py presentation.pptx --summary
    python tools/parse_ppt.py presentation.pptx --notes --format markdown
"""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    print(
        "Error: python-pptx is required. Install with: pip install python-pptx",
        file=sys.stderr,
    )
    sys.exit(1)


def parse_args():
    parser = argparse.ArgumentParser(
        description="Parse PowerPoint (.pptx) files for Copilot CLI consumption"
    )
    parser.add_argument("file", help="Path to the .pptx file")
    parser.add_argument(
        "--slides", action="store_true", help="List all slide numbers and titles"
    )
    parser.add_argument("--slide", type=int, help="Parse specific slide number (1-based)")
    parser.add_argument(
        "--format",
        choices=["json", "markdown"],
        default="json",
        help="Output format (default: json)",
    )
    parser.add_argument(
        "--images", action="store_true", help="Include image metadata"
    )
    parser.add_argument(
        "--tables", action="store_true", help="Extract tables only"
    )
    parser.add_argument(
        "--notes", action="store_true", help="Include speaker notes"
    )
    parser.add_argument(
        "--summary", action="store_true", help="Output summary statistics"
    )
    parser.add_argument("--max-slides", type=int, help="Maximum slides to output")
    return parser.parse_args()


def get_slide_title(slide):
    """Extract slide title from title placeholder."""
    if slide.shapes.title:
        return slide.shapes.title.text.strip()
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                return text[:100]
    return None


def extract_text_from_shape(shape):
    """Extract text content from a shape with formatting info."""
    if not shape.has_text_frame:
        return None

    paragraphs = []
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        para_info = {"text": text, "level": para.level}
        if para.font and para.font.bold:
            para_info["bold"] = True
        paragraphs.append(para_info)

    return paragraphs if paragraphs else None


def extract_table(shape):
    """Extract table data from a table shape."""
    if not shape.has_table:
        return None

    table = shape.table
    rows = []
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            row_data.append(cell.text.strip())
        rows.append(row_data)

    if not rows:
        return None

    headers = rows[0] if rows else []
    data_rows = rows[1:] if len(rows) > 1 else []

    return {
        "headers": headers,
        "rows": data_rows,
        "row_count": len(data_rows),
        "col_count": len(headers),
    }


def extract_image_info(shape):
    """Extract image metadata from a picture shape."""
    info = {
        "name": shape.name,
        "width": emu_to_readable(shape.width),
        "height": emu_to_readable(shape.height),
        "left": emu_to_readable(shape.left),
        "top": emu_to_readable(shape.top),
    }
    if shape.image:
        info["content_type"] = shape.image.content_type
        info["size_bytes"] = len(shape.image.blob)
        if shape.image.filename:
            info["filename"] = shape.image.filename
    return info


def emu_to_readable(emu):
    """Convert EMU to inches (rounded)."""
    if emu is None:
        return None
    return round(emu / 914400, 2)


def parse_slide(slide, slide_number, include_images=False, include_notes=False):
    """Parse a single slide into structured data."""
    slide_data = {
        "slide_number": slide_number,
        "title": get_slide_title(slide),
        "content": [],
        "tables": [],
    }

    if include_images:
        slide_data["images"] = []

    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_data = extract_table(shape)
            if table_data:
                slide_data["tables"].append(table_data)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if include_images:
                slide_data["images"].append(extract_image_info(shape))
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if child.has_text_frame:
                    paragraphs = extract_text_from_shape(child)
                    if paragraphs:
                        slide_data["content"].extend(paragraphs)
        else:
            paragraphs = extract_text_from_shape(shape)
            if paragraphs:
                slide_data["content"].extend(paragraphs)

    if include_notes and slide.has_notes_slide:
        notes_frame = slide.notes_slide.notes_text_frame
        if notes_frame and notes_frame.text.strip():
            slide_data["notes"] = notes_frame.text.strip()

    # Clean up empty lists
    if not slide_data["tables"]:
        del slide_data["tables"]
    if not slide_data["content"]:
        del slide_data["content"]
    if include_images and not slide_data.get("images"):
        slide_data.pop("images", None)

    return slide_data


def format_json(slides_data, filename):
    """Format slides data as JSON."""
    output = {
        "file": filename,
        "total_slides": len(slides_data),
        "slides": slides_data,
    }
    return json.dumps(output, ensure_ascii=False, indent=2, default=str)


def format_markdown(slides_data, filename):
    """Format slides data as Markdown."""
    lines = [f"# {filename}", ""]

    for slide in slides_data:
        num = slide["slide_number"]
        title = slide.get("title") or "(제목 없음)"
        lines.append(f"## Slide {num}: {title}")
        lines.append("")

        # Content
        if "content" in slide:
            for para in slide["content"]:
                indent = "  " * para.get("level", 0)
                text = para["text"]
                if para.get("bold"):
                    text = f"**{text}**"
                if para.get("level", 0) > 0:
                    lines.append(f"{indent}- {text}")
                else:
                    lines.append(text)
            lines.append("")

        # Tables
        if "tables" in slide:
            for table in slide["tables"]:
                headers = table["headers"]
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in table["rows"]:
                    # Pad row to match headers length
                    padded = row + [""] * (len(headers) - len(row))
                    lines.append("| " + " | ".join(padded[:len(headers)]) + " |")
                lines.append("")

        # Images
        if "images" in slide:
            lines.append("**Images:**")
            for img in slide["images"]:
                name = img.get("filename") or img.get("name", "image")
                size = img.get("size_bytes", 0)
                lines.append(f"- {name} ({size} bytes, {img['width']}\" × {img['height']}\")")
            lines.append("")

        # Notes
        if "notes" in slide:
            lines.append(f"> **Notes:** {slide['notes']}")
            lines.append("")

        lines.append("---")
        lines.append("")

    return "\n".join(lines)


def generate_summary(prs, filename):
    """Generate summary statistics about the presentation."""
    total_slides = len(prs.slides)
    total_text_shapes = 0
    total_tables = 0
    total_images = 0
    total_charts = 0
    slide_titles = []

    for i, slide in enumerate(prs.slides, 1):
        title = get_slide_title(slide)
        slide_titles.append({"slide": i, "title": title})
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                total_tables += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total_images += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
                total_charts += 1
            elif shape.has_text_frame and shape.text_frame.text.strip():
                total_text_shapes += 1

    # Slide dimensions
    width_inches = emu_to_readable(prs.slide_width)
    height_inches = emu_to_readable(prs.slide_height)

    summary = {
        "file": filename,
        "slide_dimensions": f"{width_inches}\" × {height_inches}\"",
        "total_slides": total_slides,
        "total_text_shapes": total_text_shapes,
        "total_tables": total_tables,
        "total_images": total_images,
        "total_charts": total_charts,
        "slide_titles": slide_titles,
    }
    return json.dumps(summary, ensure_ascii=False, indent=2, default=str)


def list_slides(prs, filename):
    """List all slides with their numbers and titles."""
    slides_info = []
    for i, slide in enumerate(prs.slides, 1):
        title = get_slide_title(slide)
        shape_count = len(slide.shapes)
        slides_info.append({
            "slide": i,
            "title": title,
            "shape_count": shape_count,
        })

    output = {"file": filename, "total_slides": len(prs.slides), "slides": slides_info}
    return json.dumps(output, ensure_ascii=False, indent=2)


def main():
    args = parse_args()

    filepath = Path(args.file)
    if not filepath.exists():
        print(f"Error: File not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    if filepath.suffix.lower() not in (".pptx", ".pptm"):
        print(
            f"Error: Unsupported file format '{filepath.suffix}'. "
            "Supported: .pptx, .pptm",
            file=sys.stderr,
        )
        sys.exit(1)

    prs = Presentation(str(filepath))

    if args.slides:
        print(list_slides(prs, filepath.name))
        return

    if args.summary:
        print(generate_summary(prs, filepath.name))
        return

    # Determine which slides to parse
    slides_to_parse = list(enumerate(prs.slides, 1))

    if args.slide:
        if args.slide < 1 or args.slide > len(prs.slides):
            print(
                f"Error: Slide {args.slide} out of range. "
                f"Total slides: {len(prs.slides)}",
                file=sys.stderr,
            )
            sys.exit(1)
        slides_to_parse = [(args.slide, prs.slides[args.slide - 1])]

    if args.max_slides:
        slides_to_parse = slides_to_parse[: args.max_slides]

    # Parse slides
    slides_data = []
    for slide_num, slide in slides_to_parse:
        slide_data = parse_slide(
            slide,
            slide_num,
            include_images=args.images,
            include_notes=args.notes,
        )

        # If --tables, only include slides that have tables
        if args.tables:
            if "tables" in slide_data:
                slides_data.append({
                    "slide_number": slide_data["slide_number"],
                    "title": slide_data.get("title"),
                    "tables": slide_data["tables"],
                })
        else:
            slides_data.append(slide_data)

    # Output
    if args.format == "json":
        print(format_json(slides_data, filepath.name))
    elif args.format == "markdown":
        print(format_markdown(slides_data, filepath.name))


if __name__ == "__main__":
    main()
